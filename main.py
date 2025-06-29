import sys
import re
import ctypes
from pathlib import Path
from PySide6.QtWidgets import (QApplication, QMainWindow, QVBoxLayout, 
                               QHBoxLayout, QWidget, QPushButton, 
                               QProgressBar, QLabel, QFileDialog, 
                               QMessageBox, QTextEdit, QFrame, QTableWidget, 
                               QTableWidgetItem, QHeaderView, QMenu)
from PySide6.QtCore import QThread, Signal, Qt, QUrl, QMimeData
from PySide6.QtGui import QFont, QDragEnterEvent, QDropEvent, QIcon, QColor, QAction, QClipboard
import qtawesome as qta
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.chrome.options import Options

# Libraries untuk membaca berbagai format file
try:
    from docx import Document  # untuk .docx
except ImportError:
    Document = None

try:
    import openpyxl  # untuk .xlsx
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation  # untuk .pptx
except ImportError:
    Presentation = None

try:
    import PyPDF2  # untuk .pdf
except ImportError:
    PyPDF2 = None

try:
    import xlrd  # untuk .xls
except ImportError:
    xlrd = None

try:
    import olefile  # untuk .doc dan .ppt
    import struct
except ImportError:
    olefile = None
    struct = None

try:
    import pandas as pd  # untuk .csv
except ImportError:
    pd = None

try:
    from odf.opendocument import load as odf_load  # untuk .odt, .ods, .odp
    from odf.text import P as OdfP
    from odf.table import Table as OdfTable, TableRow as OdfTableRow, TableCell as OdfTableCell
    from odf.draw import Page as OdfPage, Frame as OdfFrame
    from odf import teletype
except ImportError:
    odf_load = None
    OdfP = None
    OdfTable = None
    OdfTableRow = None
    OdfTableCell = None
    OdfPage = None
    OdfFrame = None
    teletype = None

try:
    from striprtf.striprtf import rtf_to_text  # untuk .rtf
except ImportError:
    rtf_to_text = None


class LinkOpenerWorker(QThread):
    """Worker thread untuk membuka link agar GUI tidak freeze"""
    progress_updated = Signal(int)
    status_updated = Signal(str)
    finished = Signal()
    chrome_tab_opened = Signal(str)  # Signal untuk melaporkan tab handle yang dibuka
    link_processing = Signal(int)  # Signal untuk melaporkan index link yang sedang diproses
    link_opened = Signal(int)  # Signal untuk melaporkan index link yang berhasil dibuka
    
    def __init__(self, links):
        super().__init__()
        self.links = links
        self.driver = None
        self.opened_tabs = []  # Track tab handles yang dibuka
    def run(self):
        total_links = len(self.links)
        
        try:
            # Setup Chrome driver dengan incognito mode
            self.setup_chrome_driver()
            for i, link in enumerate(self.links):
                try:                
                    # Emit signal bahwa link ini sedang diproses
                    self.link_processing.emit(i)
                    self.status_updated.emit(f"Membuka: {link}")
                    
                    # Buka link di tab baru tanpa tunggu loading
                    if i == 0:
                        # Tab pertama langsung di window utama - langsung navigate tanpa tunggu
                        self.driver.execute_script(f"window.location.href = '{link}';")
                        tab_handle = self.driver.current_window_handle
                    else:
                        # Tab berikutnya buka di window/tab baru dan langsung navigate
                        self.driver.execute_script(f"window.open('{link}');")
                        # Get handle tab yang baru dibuka
                        tab_handle = self.driver.window_handles[-1]
                    
                    # Track tab handle
                    self.opened_tabs.append(tab_handle)
                    self.chrome_tab_opened.emit(tab_handle)
                    print(f"DEBUG: Worker - Opened tab: {tab_handle} for {link}")
                    
                    # Emit signal bahwa link berhasil dibuka
                    self.link_opened.emit(i)
                    
                    # Update progress
                    progress = int((i + 1) / total_links * 100)
                    self.progress_updated.emit(progress)
                    
                    # Delay singkat untuk stabilitas (tidak tunggu loading)
                    self.msleep(200)  # Lebih cepat, cuma 200ms
                    
                except Exception as e:
                    self.status_updated.emit(f"Error membuka {link}: {str(e)}")
                    print(f"DEBUG: Worker - Error opening {link}: {e}")
        
        except Exception as e:
            self.status_updated.emit(f"Error setup Chrome driver: {str(e)}")
            print(f"DEBUG: Worker - Driver setup error: {e}")
        
        self.status_updated.emit("Selesai membuka semua link!")
        self.finished.emit()
    
    def setup_chrome_driver(self):
        """Setup Chrome driver dengan incognito mode"""
        try:
            # Path ke chromedriver.exe di folder yang sama
            chromedriver_path = Path(__file__).parent / "chromedriver.exe"
            
            if not chromedriver_path.exists():
                raise Exception("chromedriver.exe tidak ditemukan di folder aplikasi")
              # Chrome options untuk incognito mode
            chrome_options = Options()
            chrome_options.add_argument("--incognito")
            chrome_options.add_argument("--disable-web-security")
            chrome_options.add_argument("--disable-features=VizDisplayCompositor")
            
            # Hilangkan pesan "Chrome is being controlled by automated test software"
            chrome_options.add_argument("--disable-blink-features=AutomationControlled")
            chrome_options.add_experimental_option("excludeSwitches", ["enable-automation"])
            chrome_options.add_experimental_option('useAutomationExtension', False)
            
            # Setup service
            service = Service(str(chromedriver_path))
              # Buat driver instance
            self.driver = webdriver.Chrome(service=service, options=chrome_options)
            
            # Hilangkan deteksi webdriver dengan execute script
            self.driver.execute_script("Object.defineProperty(navigator, 'webdriver', {get: () => undefined})")
            
            print("DEBUG: Worker - Chrome driver setup successful")
            
        except Exception as e:
            print(f"DEBUG: Worker - Chrome driver setup failed: {e}")
            raise e
    
    def cleanup_driver(self):
        """Cleanup Chrome driver safely"""
        try:
            if self.driver:
                self.driver.quit()
                self.driver = None
                print("DEBUG: Worker - Chrome driver cleaned up")
        except Exception as e:
            print(f"DEBUG: Worker - Error cleaning up driver: {e}")


class LinkOpenerApp(QMainWindow):    
    def __init__(self):
        super().__init__()
        self.worker = None
        self.is_processing = False  # Flag untuk mencegah multiple execution
        self.source_file_path = None  # Track source file path for export
        self.opened_chrome_tabs = []  # Track Chrome tab handles yang dibuka dari app ini
        self.chrome_driver = None  # Track Chrome driver instance
        self.init_ui()
        
        # Set window always on top
        self.setWindowFlags(self.windowFlags() | Qt.WindowStaysOnTopHint)
        
        # Enable drag and drop
        self.setAcceptDrops(True)    
    def init_ui(self):
        self.setWindowTitle("Link Opener - Buka & Ekstrak Link dari File")
        self.setGeometry(100, 100, 600, 600)
        
        # Set aplikasi icon
        icon_path = Path(__file__).parent / "link_opener.ico"
        if icon_path.exists():
            self.setWindowIcon(QIcon(str(icon_path)))
        
        # Widget utama
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        # Layout utama
        layout = QVBoxLayout(central_widget)
        layout.setSpacing(20)
        layout.setContentsMargins(30, 30, 30, 30)
        
        # Icon dan Title Layout
        header_layout = QHBoxLayout()
        
        # App Icon
        icon_label = QLabel()
        icon_path = Path(__file__).parent / "link_opener.ico"
        if icon_path.exists():
            app_icon = QIcon(str(icon_path))
            # Ambil ukuran HD terbaik (256x256)
            pixmap = app_icon.pixmap(256, 256)
            # Scale ke ukuran yang sesuai untuk header (64x64)
            scaled_pixmap = pixmap.scaled(64, 64, Qt.KeepAspectRatio, Qt.SmoothTransformation)
            icon_label.setPixmap(scaled_pixmap)
        icon_label.setAlignment(Qt.AlignCenter)
        
        # Title
        title = QLabel("Link Opener")
        title.setAlignment(Qt.AlignCenter)
        title.setFont(QFont("Arial", 16, QFont.Bold))
        
        # Add ke header layout
        header_layout.addStretch()
        header_layout.addWidget(icon_label)
        header_layout.addSpacing(15)
        header_layout.addWidget(title)
        header_layout.addStretch()
        
        layout.addLayout(header_layout)        # Deskripsi
        desc = QLabel("Pilih file yang berisi campuran teks dan link.\n"
                     "Aplikasi akan menampilkan semua link yang ditemukan dalam tabel.\n"
                     "Klik 'Buka Chrome' untuk membuka semua link, atau double-click link untuk buka satu per satu.")
        desc.setAlignment(Qt.AlignCenter)
        desc.setWordWrap(True)
        layout.addWidget(desc)
        
        # Drag and Drop Area
        self.drop_frame = QFrame()
        self.drop_frame.setMinimumHeight(100)
        self.drop_frame.setStyleSheet("""
            QFrame#dropFrame {
                border: 2px dashed #aaaaaa;
                border-radius: 10px;
            }
        """)
        self.drop_frame.setObjectName("dropFrame")
        drop_layout = QVBoxLayout(self.drop_frame)
        
        # Main drag & drop label
        drop_label = QLabel("Drag & Drop file di sini\natau klik tombol di bawah")
        drop_label.setAlignment(Qt.AlignCenter)
        drop_label.setFont(QFont("Arial", 10))
        drop_layout.addWidget(drop_label)
        
        # Supported formats label with smaller font
        formats_label = QLabel("Mendukung format: TXT, DOC, DOCX, XLS, XLSX, PPT, PPTX, PDF, CSV, RTF, ODT, ODS, ODP")
        formats_label.setAlignment(Qt.AlignCenter)
        formats_label.setFont(QFont("Arial", 8))
        formats_label.setStyleSheet("color: #666666; margin-top: 5px;")
        formats_label.setWordWrap(True)
        drop_layout.addWidget(formats_label)
        
        layout.addWidget(self.drop_frame)
        
        # Layout horizontal untuk tombol-tombol di bawah DND
        buttons_layout = QHBoxLayout()
          # Tombol buka file
        self.open_button = QPushButton("Pilih File")
        self.open_button.setIcon(qta.icon('fa5s.folder-open', color='#2196F3'))
        self.open_button.setMinimumHeight(50)
        self.open_button.setFont(QFont("Arial", 12))
        self.open_button.setStyleSheet("""
            QPushButton {
                border: 2px solid #2196F3;
                border-radius: 8px;
                padding: 8px 16px;
                background-color: transparent;
            }
            QPushButton:hover {
                border: 2px solid #1976D2;
                background-color: rgba(33, 150, 243, 0.05);
            }
            QPushButton:pressed {
                border: 2px solid #0D47A1;
                background-color: rgba(33, 150, 243, 0.1);
            }
        """)
        self.open_button.clicked.connect(self.open_file)
        buttons_layout.addWidget(self.open_button)
          # Tombol buka link
        self.open_links_button = QPushButton("Buka Chrome")
        self.open_links_button.setIcon(qta.icon('fa5s.rocket', color='#4CAF50'))
        self.open_links_button.setMinimumHeight(50)
        self.open_links_button.setFont(QFont("Arial", 12))
        self.open_links_button.setStyleSheet("""
            QPushButton {
                border: 2px solid #4CAF50;
                border-radius: 8px;
                padding: 8px 16px;
                background-color: transparent;
            }
            QPushButton:hover {
                border: 2px solid #388E3C;
                background-color: rgba(76, 175, 80, 0.05);
            }
            QPushButton:pressed {
                border: 2px solid #2E7D32;
                background-color: rgba(76, 175, 80, 0.1);
            }
            QPushButton:disabled {
                border: 2px solid #CCCCCC;
                color: #999999;
            }
        """)
        self.open_links_button.setVisible(False)
        self.open_links_button.clicked.connect(self.open_links)
        buttons_layout.addWidget(self.open_links_button)
          # Tombol export links
        self.export_button = QPushButton("Export Links")
        self.export_button.setIcon(qta.icon('fa5s.file-export', color='#FF9800'))
        self.export_button.setMinimumHeight(50)
        self.export_button.setFont(QFont("Arial", 12))
        self.export_button.setStyleSheet("""
            QPushButton {
                border: 2px solid #FF9800;
                border-radius: 8px;
                padding: 8px 16px;
                background-color: transparent;
            }
            QPushButton:hover {
                border: 2px solid #F57C00;
                background-color: rgba(255, 152, 0, 0.05);
            }
            QPushButton:pressed {
                border: 2px solid #E65100;
                background-color: rgba(255, 152, 0, 0.1);
            }
        """)
        self.export_button.setVisible(False)
        self.export_button.clicked.connect(self.export_links)
        buttons_layout.addWidget(self.export_button)
        
        # Tombol tutup tab
        self.close_tabs_button = QPushButton("Tutup Tab")
        self.close_tabs_button.setIcon(qta.icon('fa5s.times-circle', color='#F44336'))
        self.close_tabs_button.setMinimumHeight(50)
        self.close_tabs_button.setFont(QFont("Arial", 12))
        self.close_tabs_button.setStyleSheet("""
            QPushButton {
                border: 2px solid #F44336;
                border-radius: 8px;
                padding: 8px 16px;
                background-color: transparent;
            }
            QPushButton:hover {
                border: 2px solid #D32F2F;
                background-color: rgba(244, 67, 54, 0.05);
            }
            QPushButton:pressed {
                border: 2px solid #B71C1C;
                background-color: rgba(244, 67, 54, 0.1);
            }
            QPushButton:disabled {
                border: 2px solid #CCCCCC;
                color: #999999;
            }
        """)
        self.close_tabs_button.setVisible(False)
        self.close_tabs_button.clicked.connect(self.close_chrome_tabs)
        buttons_layout.addWidget(self.close_tabs_button)
        
        layout.addLayout(buttons_layout)
        
        # Status label
        self.status_label = QLabel("")
        self.status_label.setAlignment(Qt.AlignLeft)
        self.status_label.setWordWrap(True)
        layout.addWidget(self.status_label)        # Table widget untuk menampilkan link yang ditemukan
        self.links_table = QTableWidget()
        self.links_table.setColumnCount(1)
        self.links_table.setHorizontalHeaderLabels(["Link"])
        self.links_table.setVisible(False)
        self.links_table.setSelectionBehavior(QTableWidget.SelectRows)
        self.links_table.setEditTriggers(QTableWidget.NoEditTriggers)
        
        # Connect double-click event untuk buka link individual
        self.links_table.itemDoubleClicked.connect(self.open_single_link)
        
        # Enable context menu untuk table
        self.links_table.setContextMenuPolicy(Qt.CustomContextMenu)
        self.links_table.customContextMenuRequested.connect(self.show_context_menu)
        
        # Set column widths
        header = self.links_table.horizontalHeader()
        header.setSectionResizeMode(0, QHeaderView.Stretch)  # Link column stretches
        
        layout.addWidget(self.links_table)
        
        # Progress bar (moved below table)
        self.progress_bar = QProgressBar()
        self.progress_bar.setVisible(False)
        self.progress_bar.setMinimumHeight(25)
        self.progress_bar.setFormat("Sedang membuka link...")  # Custom text instead of percentage
        layout.addWidget(self.progress_bar)
        layout.addStretch()
        
        self.found_links = []
    
    def dragEnterEvent(self, event: QDragEnterEvent):
        """Handle drag enter event"""
        if event.mimeData().hasUrls():
            urls = event.mimeData().urls()
            if len(urls) == 1:
                file_path = urls[0].toLocalFile().lower()
                supported_extensions = ['.txt', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.pdf', '.csv', '.rtf', '.odt', '.ods', '.odp']
                if any(file_path.endswith(ext) for ext in supported_extensions):
                    # Set hover style when dragging valid file - only border color change
                    self.drop_frame.setStyleSheet("""
                        QFrame#dropFrame {
                            border: 2px dashed #4CAF50;
                            border-radius: 10px;
                        }
                    """)
                    event.acceptProposedAction()
                else:
                    event.ignore()
            else:
                event.ignore()
        else:
            event.ignore()
    
    def dragLeaveEvent(self, event):
        """Handle drag leave event"""
        # Reset to normal style when drag leaves
        self.drop_frame.setStyleSheet("""
            QFrame#dropFrame {
                border: 2px dashed #aaaaaa;
                border-radius: 10px;
            }
        """)
        event.accept()    
    def dropEvent(self, event: QDropEvent):
        """Handle drop event"""
        # Reset to normal style first
        self.drop_frame.setStyleSheet("""
            QFrame#dropFrame {
                border: 2px dashed #aaaaaa;
                border-radius: 10px;
            }
        """)
        
        if event.mimeData().hasUrls():
            urls = event.mimeData().urls()
            if len(urls) == 1:
                file_path = urls[0].toLocalFile()
                file_path_lower = file_path.lower()
                supported_extensions = ['.txt', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.pdf', '.csv', '.rtf', '.odt', '.ods', '.odp']
                if any(file_path_lower.endswith(ext) for ext in supported_extensions):
                    self.load_and_extract_links(file_path)
                    event.acceptProposedAction()                
                else:
                    QMessageBox.warning(self, "Peringatan", "Format file tidak didukung!\nHanya mendukung: TXT, DOC, DOCX, XLS, XLSX, PPT, PPTX, PDF, CSV, RTF, ODT, ODS, ODP")
            else:
                QMessageBox.warning(self, "Peringatan", "Hanya bisa drop satu file!")
        else:
            event.ignore()
    def open_file(self):
        """Buka dialog untuk memilih file"""
        file_path, _ = QFileDialog.getOpenFileName(
            self, 
            "Pilih File", 
            "", 
            "All Supported (*.txt *.doc *.docx *.xls *.xlsx *.ppt *.pptx *.pdf *.csv *.rtf *.odt *.ods *.odp);;"
            "Text Files (*.txt);;"
            "Word Documents (*.doc *.docx);;"
            "Excel Files (*.xls *.xlsx);;"
            "PowerPoint Files (*.ppt *.pptx);;"
            "PDF Files (*.pdf);;"
            "CSV Files (*.csv);;"
            "RTF Files (*.rtf);;"
            "OpenDocument Files (*.odt *.ods *.odp);;"
            "All Files (*)"
        )
        
        if file_path:
            self.load_and_extract_links(file_path)
    def extract_text_from_file(self, file_path):
        """Ekstrak teks dari berbagai format file"""
        file_extension = Path(file_path).suffix.lower()
        
        try:
            if file_extension == '.txt':
                return self.extract_text_from_txt(file_path)
            elif file_extension == '.docx':
                return self.extract_text_from_docx(file_path)
            elif file_extension == '.doc':
                return self.extract_text_from_doc(file_path)
            elif file_extension in ['.xlsx', '.xls']:
                return self.extract_text_from_excel(file_path)
            elif file_extension == '.pptx':
                return self.extract_text_from_pptx(file_path)
            elif file_extension == '.ppt':
                return self.extract_text_from_ppt(file_path)
            elif file_extension == '.pdf':
                return self.extract_text_from_pdf(file_path)
            elif file_extension == '.csv':
                return self.extract_text_from_csv(file_path)
            elif file_extension == '.rtf':
                return self.extract_text_from_rtf(file_path)
            elif file_extension == '.odt':
                return self.extract_text_from_odt(file_path)
            elif file_extension == '.ods':
                return self.extract_text_from_ods(file_path)
            elif file_extension == '.odp':
                return self.extract_text_from_odp(file_path)
            else:
                raise Exception(f"Format file {file_extension} tidak didukung")
                
        except Exception as e:
            raise Exception(f"Gagal membaca file {file_extension}: {str(e)}")
    
    def extract_text_from_txt(self, file_path):
        """Ekstrak teks dari file TXT"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            return file.read()
    
    def extract_text_from_docx(self, file_path):
        """Ekstrak teks dari file DOCX termasuk hyperlink"""
        if Document is None:
            raise Exception("Library python-docx tidak terinstall. Install dengan: pip install python-docx")
        
        doc = Document(file_path)
        text = []
        links = []
        
        # Ekstrak teks normal dan hyperlink
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
            
            # Ekstrak hyperlink dari paragraph
            for run in paragraph.runs:
                if run.element.tag.endswith('hyperlink') or run._element.getparent().tag.endswith('hyperlink'):
                    # Coba ambil URL dari hyperlink
                    hyperlink_element = run._element.getparent()
                    if hyperlink_element.tag.endswith('hyperlink'):
                        rId = hyperlink_element.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                        if rId:
                            try:
                                relationship = doc.part.rels[rId]
                                if relationship.target_ref.startswith('http'):
                                    links.append(relationship.target_ref)
                            except:
                                pass
        
        # Gabungkan teks dan link yang ditemukan
        all_text = '\n'.join(text)
        if links:
            all_text += '\n' + '\n'.join(links)
        
        return all_text
    
    def extract_text_from_doc(self, file_path):
        """Ekstrak teks dari file DOC (format lama)"""
        if olefile is None:
            raise Exception("Library olefile tidak terinstall. Install dengan: pip install olefile")
        
        # Untuk file .doc, kita coba baca sebagai binary dan cari teks
        try:
            with open(file_path, 'rb') as file:
                content = file.read()
                # Decode dengan berbagai encoding
                text = content.decode('utf-8', errors='ignore')
                # Filter karakter yang tidak dapat dibaca
                text = ''.join(char for char in text if ord(char) >= 32 or char in '\n\t')
                return text
        except Exception:
            raise Exception("Gagal membaca file DOC. Coba convert ke DOCX dulu.")
    
    def extract_text_from_excel(self, file_path):
        """Ekstrak teks dari file Excel (XLS/XLSX) termasuk hyperlink"""
        file_extension = Path(file_path).suffix.lower()
        text = []
        
        if file_extension == '.xlsx':
            if openpyxl is None:
                raise Exception("Library openpyxl tidak terinstall. Install dengan: pip install openpyxl")
            
            workbook = openpyxl.load_workbook(file_path, data_only=False)  # Keep formulas untuk hyperlink
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                for row in sheet.iter_rows():
                    for cell in row:
                        # Tambahkan nilai cell
                        if cell.value is not None:
                            text.append(str(cell.value))
                        
                        # Cek hyperlink
                        if cell.hyperlink is not None:
                            if hasattr(cell.hyperlink, 'target') and cell.hyperlink.target:
                                if cell.hyperlink.target.startswith('http'):
                                    text.append(cell.hyperlink.target)
            workbook.close()
            
        elif file_extension == '.xls':
            if xlrd is None:
                raise Exception("Library xlrd tidak terinstall. Install dengan: pip install xlrd")
            
            workbook = xlrd.open_workbook(file_path, formatting_info=True)
            for sheet in workbook.sheets():
                for row in range(sheet.nrows):
                    for col in range(sheet.ncols):
                        cell_value = sheet.cell_value(row, col)
                        if cell_value:
                            text.append(str(cell_value))
                        
                        # Coba ekstrak hyperlink dari XLS (lebih kompleks)
                        try:
                            cell_obj = sheet.cell(row, col)
                            if hasattr(cell_obj, 'ctype') and cell_obj.ctype == xlrd.XL_CELL_TEXT:
                                # Untuk XLS, hyperlink biasanya tersimpan sebagai text yang dimulai dengan http
                                cell_text = str(cell_value)
                                if cell_text.startswith('http'):
                                    text.append(cell_text)
                        except:
                            pass
        
        return '\n'.join(text)
    
    def extract_text_from_pptx(self, file_path):
        """Ekstrak teks dari file PPTX termasuk hyperlink"""
        if Presentation is None:
            raise Exception("Library python-pptx tidak terinstall. Install dengan: pip install python-pptx")
        
        prs = Presentation(file_path)
        text = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
                
                # Cek hyperlink di shape
                if hasattr(shape, "click_action") and shape.click_action.hyperlink:
                    hyperlink = shape.click_action.hyperlink
                    if hasattr(hyperlink, 'address') and hyperlink.address:
                        if hyperlink.address.startswith('http'):
                            text.append(hyperlink.address)
                
                # Cek hyperlink di text runs (untuk text yang ada hyperlink-nya)
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if hasattr(run, "hyperlink") and run.hyperlink:
                                if hasattr(run.hyperlink, 'address') and run.hyperlink.address:
                                    if run.hyperlink.address.startswith('http'):
                                        text.append(run.hyperlink.address)
        
        return '\n'.join(text)
    
    def extract_text_from_pdf(self, file_path):
        """Ekstrak teks dari file PDF termasuk link/annotation"""
        if PyPDF2 is None:
            raise Exception("Library PyPDF2 tidak terinstall. Install dengan: pip install PyPDF2")
        
        text = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    # Ekstrak teks normal
                    page_text = page.extract_text()
                    if page_text:
                        text.append(page_text)
                    
                    # Coba ekstrak link dari annotations
                    if hasattr(page, 'annotations') and page.annotations:
                        for annotation in page.annotations:
                            if annotation.get_object():
                                annotation_obj = annotation.get_object()
                                if '/A' in annotation_obj:
                                    action = annotation_obj['/A']
                                    if '/URI' in action:
                                        uri = action['/URI']
                                        if isinstance(uri, str) and uri.startswith('http'):
                                            text.append(uri)
        except Exception as e:
            raise Exception(f"Gagal membaca PDF: {str(e)}")
        
        return '\n'.join(text)
    
    def extract_text_from_ppt(self, file_path):
        """Ekstrak teks dari file PPT (format lama)"""
        if olefile is None:
            raise Exception("Library olefile tidak terinstall. Install dengan: pip install olefile")
        
        # Untuk file .ppt, coba baca sebagai binary dan cari teks
        try:
            with open(file_path, 'rb') as file:
                content = file.read()
                # Decode dengan berbagai encoding
                text = content.decode('utf-8', errors='ignore')
                # Filter karakter yang tidak dapat dibaca
                text = ''.join(char for char in text if ord(char) >= 32 or char in '\n\t')
                return text
        except Exception:
            raise Exception("Gagal membaca file PPT. Coba convert ke PPTX dulu.")
    
    def extract_text_from_csv(self, file_path):
        """Ekstrak teks dari file CSV"""
        if pd is None:
            raise Exception("Library pandas tidak terinstall. Install dengan: pip install pandas")
        
        text = []
        try:
            # Coba baca dengan encoding UTF-8 dulu
            df = pd.read_csv(file_path, encoding='utf-8')
        except UnicodeDecodeError:
            try:
                # Kalau gagal, coba dengan encoding latin-1
                df = pd.read_csv(file_path, encoding='latin-1')
            except:
                # Terakhir coba dengan encoding cp1252 (Windows)
                df = pd.read_csv(file_path, encoding='cp1252')
        
        # Ekstrak semua nilai dari DataFrame
        for column in df.columns:
            # Tambahkan nama kolom
            text.append(str(column))
            # Tambahkan nilai kolom (drop NaN values)
            values = df[column].dropna().astype(str).tolist()
            text.extend(values)
        
        return '\n'.join(text)
    
    def extract_text_from_rtf(self, file_path):
        """Ekstrak teks dari file RTF"""
        if rtf_to_text is None:
            raise Exception("Library striprtf tidak terinstall. Install dengan: pip install striprtf")
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                rtf_content = file.read()
            
            # Convert RTF ke plain text
            plain_text = rtf_to_text(rtf_content)
            return plain_text
        except UnicodeDecodeError:
            # Coba dengan encoding lain
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    rtf_content = file.read()
                plain_text = rtf_to_text(rtf_content)
                return plain_text
            except:
                raise Exception("Gagal membaca file RTF dengan encoding yang didukung")
    
    def extract_text_from_odt(self, file_path):
        """Ekstrak teks dari file ODT (OpenDocument Text)"""
        if odf_load is None or teletype is None:
            raise Exception("Library odfpy tidak terinstall. Install dengan: pip install odfpy")
        
        try:
            doc = odf_load(file_path)
            text = []
            
            # Ekstrak semua paragraf teks
            for paragraph in doc.getElementsByType(OdfP):
                para_text = teletype.extractText(paragraph)
                if para_text.strip():
                    text.append(para_text)
            
            return '\n'.join(text)
        except Exception as e:
            raise Exception(f"Gagal membaca file ODT: {str(e)}")
    
    def extract_text_from_ods(self, file_path):
        """Ekstrak teks dari file ODS (OpenDocument Spreadsheet)"""
        if odf_load is None or teletype is None:
            raise Exception("Library odfpy tidak terinstall. Install dengan: pip install odfpy")
        
        try:
            doc = odf_load(file_path)
            text = []
            
            # Ekstrak teks dari semua tabel
            for table in doc.getElementsByType(OdfTable):
                for row in table.getElementsByType(OdfTableRow):
                    for cell in row.getElementsByType(OdfTableCell):
                        cell_text = teletype.extractText(cell)
                        if cell_text.strip():
                            text.append(cell_text)
            
            return '\n'.join(text)
        except Exception as e:
            raise Exception(f"Gagal membaca file ODS: {str(e)}")
    
    def extract_text_from_odp(self, file_path):
        """Ekstrak teks dari file ODP (OpenDocument Presentation)"""
        if odf_load is None or teletype is None:
            raise Exception("Library odfpy tidak terinstall. Install dengan: pip install odfpy")
        
        try:
            doc = odf_load(file_path)
            text = []
            
            # Ekstrak teks dari semua halaman presentasi
            for page in doc.getElementsByType(OdfPage):
                # Ekstrak teks dari frame di halaman
                for frame in page.getElementsByType(OdfFrame):
                    frame_text = teletype.extractText(frame)
                    if frame_text.strip():
                        text.append(frame_text)
                
                # Ekstrak paragraf langsung dari halaman
                for paragraph in page.getElementsByType(OdfP):
                    para_text = teletype.extractText(paragraph)
                    if para_text.strip():
                        text.append(para_text)
            
            return '\n'.join(text)
        except Exception as e:
            raise Exception(f"Gagal membaca file ODP: {str(e)}")
    
    def load_and_extract_links(self, file_path):
        """Load file dan ekstrak link dengan metode yang diperbaiki"""
        try:
            # Store source file path for export functionality
            self.source_file_path = file_path
            
            # Ekstrak teks berdasarkan format file
            content = self.extract_text_from_file(file_path)
            
            # Ambil nama file tanpa path
            file_name = Path(file_path).name
            
            # Regex yang diperbaiki untuk mencari link http/https
            # Pattern yang lebih fleksibel untuk menangkap berbagai format link
            link_patterns = [
                r'https?://[^\s<>"\'`\[\]{}|\\^]+',  # Standard HTTP links
                r'www\.[^\s<>"\'`\[\]{}|\\^]+\.[a-zA-Z]{2,}[^\s<>"\'`\[\]{}|\\^]*',  # www links
                r'[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}/[^\s<>"\'`\[\]{}|\\^]*',  # domain/path links
            ]
            
            links = []
            for pattern in link_patterns:
                found_links = re.findall(pattern, content, re.IGNORECASE)
                links.extend(found_links)
            
            # Bersihkan dan validasi link
            self.found_links = []
            for link in links:
                # Hapus karakter yang tidak diinginkan di akhir
                cleaned_link = re.sub(r'[.,!?;:)}\]]+$', '', link.strip())
                
                # Pastikan link memiliki protocol
                if cleaned_link:
                    if not cleaned_link.startswith(('http://', 'https://')):
                        if cleaned_link.startswith('www.'):
                            cleaned_link = 'https://' + cleaned_link
                        elif '.' in cleaned_link and not cleaned_link.startswith('mailto:'):
                            # Cek apakah ini seperti domain
                            domain_pattern = r'^[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}'
                            if re.match(domain_pattern, cleaned_link):
                                cleaned_link = 'https://' + cleaned_link
                    
                    # Validasi final
                    if cleaned_link.startswith(('http://', 'https://')):
                        # Pastikan domain valid
                        try:
                            from urllib.parse import urlparse
                            parsed = urlparse(cleaned_link)
                            if parsed.netloc and '.' in parsed.netloc:
                                self.found_links.append(cleaned_link)
                        except:
                            # Jika parsing gagal, tetap tambahkan jika format basic benar
                            if re.match(r'https?://[^.]+\..+', cleaned_link):
                                self.found_links.append(cleaned_link)
            
            # Hapus duplikat
            self.found_links = list(dict.fromkeys(self.found_links))  # Preserves order
            
            if self.found_links:
                self.status_label.setText(f"Ditemukan {len(self.found_links)} link unik dari {file_name}:")
                  # Populate table dengan link
                self.links_table.setRowCount(len(self.found_links))               
                for i, link in enumerate(self.found_links):
                    # Kolom Link
                    link_item = QTableWidgetItem(link)
                    self.links_table.setItem(i, 0, link_item)
                
                self.links_table.setVisible(True)
                self.open_links_button.setVisible(True)
                self.export_button.setVisible(True)
                self.close_tabs_button.setVisible(True)
                # Disable tombol tutup tab karena belum ada tab yang terbuka
                self.close_tabs_button.setEnabled(False)
            else:
                self.status_label.setText(f"Tidak ada link yang ditemukan dalam {file_name}.")                
                self.links_table.setVisible(False)
                self.open_links_button.setVisible(False)
                self.export_button.setVisible(False)
                self.close_tabs_button.setVisible(False)
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Gagal baca file: {str(e)}")
    def open_links(self):
        """Mulai proses membuka link"""
        # Cegah multiple execution dengan disable button
        if self.is_processing:
            return            
        if not self.found_links:
            QMessageBox.warning(self, "Peringatan", "Tidak ada link untuk dibuka!")
            return
        
        # Set flag processing dan disable button
        self.is_processing = True
        self.open_links_button.setEnabled(False)
        
        # Langsung mulai tanpa konfirmasi
        self.start_opening_links()    
    def start_opening_links(self):
        """Mulai worker thread untuk membuka link"""
        self.open_links_button.setEnabled(False)
        self.open_button.setEnabled(False)
        self.progress_bar.setVisible(True)
        self.progress_bar.setValue(0)
          # Reset daftar Chrome tabs yang terbuka
        self.opened_chrome_tabs = []
        self.chrome_driver = None  # Track Chrome driver instance        # Buat dan jalankan worker thread
        self.worker = LinkOpenerWorker(self.found_links)
        self.worker.progress_updated.connect(self.progress_bar.setValue)
        self.worker.status_updated.connect(self.update_progress_text)
        self.worker.finished.connect(self.on_finished)
        self.worker.chrome_tab_opened.connect(self.track_chrome_tab)
        self.worker.link_processing.connect(self.mark_link_processing)
        self.worker.link_opened.connect(self.mark_link_opened)
        self.worker.start()    
    def track_chrome_tab(self, tab_handle):
        """Track Chrome tab handle yang dibuka dari app ini"""
        self.opened_chrome_tabs.append(tab_handle)
        print(f"DEBUG: Tracking Chrome tab: {tab_handle}, Total tracked: {len(self.opened_chrome_tabs)}")
        # Update status label dengan info tab
        self.status_label.setText(f"Melacak {len(self.opened_chrome_tabs)} tab Chrome...")
          # Store driver reference untuk close tabs nanti
        if self.worker and self.worker.driver:
            self.chrome_driver = self.worker.driver
    
    def mark_link_processing(self, index):
        """Tandai link yang sedang diproses dengan warna kuning transparan"""
        try:
            # Set background kuning transparan untuk row yang sedang diproses menggunakan rgba
            # rgba(246, 191, 17, 0.1) - warna kuning/amber transparan
            for column in range(self.links_table.columnCount()):
                item = self.links_table.item(index, column)
                if item:
                    # Gunakan QColor dengan nilai rgba yang dikonversi
                    rgba_color = QColor(246, 191, 17, 25)  # rgba(246, 191, 17, 0.1)
                    item.setBackground(rgba_color)
            print(f"DEBUG: Marked link {index} as processing (rgba yellow: 246, 191, 17, 0.1)")
        except Exception as e:
            print(f"DEBUG: Error marking link {index} as processing: {e}")
    
    def mark_link_opened(self, index):
        """Tandai link yang berhasil dibuka dengan warna hijau transparan"""
        try:
            # Set background hijau transparan untuk row yang berhasil dibuka menggunakan rgba
            # rgba(100, 195, 0, 0.1) - warna hijau transparan
            rgba_color = QColor(100, 195, 0, 25)  # rgba(100, 195, 0, 0.1)
            for column in range(self.links_table.columnCount()):
                item = self.links_table.item(index, column)
                if item:
                    item.setBackground(rgba_color)
            print(f"DEBUG: Marked link {index} as opened (rgba green: 100, 195, 0, 0.1)")
        except Exception as e:
            print(f"DEBUG: Error marking link {index} as opened: {e}")
    def update_progress_text(self, text):
        """Update progress bar text dengan link yang sedang dibuka"""
        if text.startswith("Membuka:"):
            link = text.replace("Membuka: ", "")
            # Potong link jika terlalu panjang
            if len(link) > 60:
                link = link[:57] + "..."
            self.progress_bar.setFormat(f"Membuka: {link}")
        else:
            self.progress_bar.setFormat(text)    
    def on_finished(self):
        """Callback ketika selesai membuka semua link"""
        print("DEBUG: on_finished called")
        self.progress_bar.setVisible(False)
        self.open_links_button.setEnabled(True)
        self.open_button.setEnabled(True)
        
        # Reset flag processing
        self.is_processing = False
        
        # Store Chrome driver reference untuk close tabs nanti (JANGAN cleanup di sini)
        if self.worker and self.worker.driver:
            self.chrome_driver = self.worker.driver
            print(f"DEBUG: Chrome driver stored for later cleanup, tabs tracked: {len(self.opened_chrome_tabs)}")
        
        # Safely cleanup worker thread TANPA cleanup driver
        if self.worker:
            try:                 # Disconnect signals to prevent recursive calls
                self.worker.progress_updated.disconnect()
                self.worker.status_updated.disconnect()
                self.worker.finished.disconnect()                
                self.worker.chrome_tab_opened.disconnect()
                self.worker.link_processing.disconnect()
                self.worker.link_opened.disconnect()
                
                # Let thread finish naturally
                self.worker.quit()
                self.worker.wait(5000)  # Wait max 5 seconds
                
            except Exception as e:
                print(f"DEBUG: Error during worker cleanup: {e}")
            finally:
                self.worker = None
          # Ubah status tanpa dialog konfirmasi
        self.status_label.setText(f"Selesai membuka {len(self.found_links)} link di Chrome incognito!")
        
        # Update informasi tentang tab yang dibuka dan enable tombol tutup tab
        if self.opened_chrome_tabs:
            self.status_label.setText(f"Selesai membuka {len(self.found_links)} link di Chrome incognito! "
                                    f"({len(self.opened_chrome_tabs)} tab Chrome dibuka)")
            self.close_tabs_button.setEnabled(True)  # Enable tombol tutup tab
        else:
            self.close_tabs_button.setEnabled(False)  # Disable jika tidak ada tab
    
    def export_links(self):
        """Export link ke file _links.txt"""
        if not self.found_links:
            QMessageBox.warning(self, "Peringatan", "Tidak ada link untuk diekspor!")
            return
        
        if not self.source_file_path:
            QMessageBox.warning(self, "Peringatan", "Path file sumber tidak ditemukan!")
            return
        
        # Buat nama file berdasarkan source file dengan suffix _links
        source_path = Path(self.source_file_path)
        export_file_path = source_path.parent / f"{source_path.stem}_links.txt"
        
        # Cek apakah file sudah ada - hanya tampilkan dialog jika file exists
        if export_file_path.exists():
            reply = QMessageBox.question(
                self,
                "File Sudah Ada",
                f"File {export_file_path.name} sudah ada.\nIngin menimpa file yang lama?",
                QMessageBox.Yes | QMessageBox.No,
                QMessageBox.No
            )
            
            if reply == QMessageBox.No:
                return
        
        try:
            with open(export_file_path, 'w', encoding='utf-8') as file:                
                for i, link in enumerate(self.found_links, 1):
                    file.write(f"{i}. {link}\n")
            QMessageBox.information(
                self, 
                "Berhasil!", 
                f"Berhasil mengekspor {len(self.found_links)} link ke:\n{export_file_path}"
            )
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Gagal ekspor file: {str(e)}")
    
    def close_chrome_tabs(self):
        """Tutup hanya tab Chrome yang dibuka dari aplikasi ini"""
        print(f"DEBUG: close_chrome_tabs called. Current tracked tabs: {self.opened_chrome_tabs}")
        
        if not self.opened_chrome_tabs:
            self.status_label.setText("Tidak ada tab Chrome dari app ini yang perlu ditutup.")
            return
        
        if not self.chrome_driver:
            self.status_label.setText("Chrome driver tidak tersedia!")
            return
        
        try:
            closed_count = 0
            failed_count = 0
              # Tutup tab satu per satu
            for tab_handle in self.opened_chrome_tabs[:]:  # Copy list untuk iterasi aman
                try:
                    # Switch ke tab yang mau ditutup
                    self.chrome_driver.switch_to.window(tab_handle)
                    
                    # Tutup tab
                    self.chrome_driver.close()
                    
                    # Hapus dari tracking list
                    self.opened_chrome_tabs.remove(tab_handle)
                    closed_count += 1
                    print(f"DEBUG: Successfully closed tab: {tab_handle}")
                    
                except Exception as e:
                    print(f"DEBUG: Error closing tab {tab_handle}: {e}")
                    failed_count += 1
                    # Tetap hapus dari list meskipun gagal tutup
                    if tab_handle in self.opened_chrome_tabs:
                        self.opened_chrome_tabs.remove(tab_handle)
            # Setelah tutup semua tabs, cleanup driver
            if self.chrome_driver:
                try:
                    self.chrome_driver.quit()
                    self.chrome_driver = None
                    print("DEBUG: Chrome driver cleaned up after closing tabs")
                except Exception as e:
                    print(f"DEBUG: Error cleaning up driver: {e}")
              # Reset styling di table rows setelah tutup tabs
            self.reset_table_styling()
            
            # Update status label setelah tutup tabs
            if closed_count > 0:
                self.status_label.setText(f"Selesai menutup {closed_count} tab Chrome. Semua tab dari app ini sudah ditutup.")
            else:
                self.status_label.setText("Semua tab Chrome dari app ini sudah ditutup.")
              # Disable tombol tutup tab karena sudah tidak ada tab yang terbuka
            self.close_tabs_button.setEnabled(False)
                
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Gagal tutup tab Chrome: {str(e)}")
    def reset_table_styling(self):
        """Reset semua background color di table rows ke style original"""
        try:
            for row in range(self.links_table.rowCount()):
                for column in range(self.links_table.columnCount()):
                    item = self.links_table.item(row, column)
                    if item:
                        # Reset ke background default table dengan membuat item baru
                        # Simpan text content
                        original_text = item.text()
                        # Buat item baru dengan styling default
                        new_item = QTableWidgetItem(original_text)
                        # Replace item dengan yang baru (akan menggunakan style default)
                        self.links_table.setItem(row, column, new_item)
            print("DEBUG: Reset all table row styling to original default")
        except Exception as e:
            print(f"DEBUG: Error resetting table styling: {e}")
    
    def show_context_menu(self, position):
        """Tampilkan context menu saat klik kanan pada tabel"""
        item = self.links_table.itemAt(position)
        if not item:
            return
        
        # Buat context menu
        context_menu = QMenu(self)
        
        # Action untuk copy link
        copy_action = QAction("Copy Link", self)
        copy_action.setIcon(qta.icon('fa5s.copy', color='#666666'))
        copy_action.triggered.connect(lambda: self.copy_link_to_clipboard(item))
        context_menu.addAction(copy_action)
        
        # Action untuk open link
        open_action = QAction("Open Link", self)
        open_action.setIcon(qta.icon('fa5s.external-link-alt', color='#4CAF50'))
        open_action.triggered.connect(lambda: self.open_single_link(item))
        context_menu.addAction(open_action)
        
        # Tampilkan menu di posisi kursor
        context_menu.exec(self.links_table.mapToGlobal(position))
    
    def copy_link_to_clipboard(self, item):
        """Copy link ke clipboard"""
        try:
            if not item:
                return
            
            link = item.text()
            if link:
                # Get clipboard
                clipboard = QApplication.clipboard()
                clipboard.setText(link)
                
                # Show temporary status message
                original_text = self.status_label.text()
                self.status_label.setText(f"Link berhasil disalin ke clipboard: {link[:50]}...")
                
                # Reset status after 2 seconds
                from PySide6.QtCore import QTimer
                QTimer.singleShot(2000, lambda: self.status_label.setText(original_text))
                
        except Exception as e:
            QMessageBox.warning(self, "Peringatan", f"Gagal copy link: {str(e)}")
    
    def open_single_link(self, item):
        """Buka link individual saat double-click pada item tabel"""
        try:
            if not item:
                return
            
            # Ambil link dari item yang diklik
            link = item.text()
            if not link or not (link.startswith('http://') or link.startswith('https://')):
                QMessageBox.warning(self, "Peringatan", "Link tidak valid!")
                return
            
            # Cari index link di table untuk visual feedback
            link_index = -1
            for row in range(self.links_table.rowCount()):
                table_item = self.links_table.item(row, 0)
                if table_item and table_item.text() == link:
                    link_index = row
                    break
            
            # Tandai link sebagai processing (kuning)
            if link_index >= 0:
                self.mark_link_processing(link_index)
            
            # Validasi Chrome driver session atau buat baru jika perlu
            if not self.is_chrome_driver_valid():
                try:
                    self.setup_chrome_driver_for_single_link()
                except Exception as e:
                    QMessageBox.critical(self, "Error", f"Gagal setup Chrome driver: {str(e)}")
                    return
            
            # Buka link di tab baru atau window utama
            try:
                if len(self.opened_chrome_tabs) == 0:
                    # Kalau belum ada tab, langsung navigate ke link
                    self.chrome_driver.get(link)
                    tab_handle = self.chrome_driver.current_window_handle
                else:
                    # Kalau sudah ada tab, buka di tab baru
                    self.chrome_driver.execute_script(f"window.open('{link}');")
                    tab_handle = self.chrome_driver.window_handles[-1]
                  # Track tab yang baru dibuka
                self.opened_chrome_tabs.append(tab_handle)
                print(f"DEBUG: Single link opened: {link}, tab handle: {tab_handle}")
                
                # Tandai link sebagai berhasil dibuka (hijau)
                if link_index >= 0:
                    self.mark_link_opened(link_index)
                
                # Update status dan enable tombol tutup tab
                self.status_label.setText(f"Link dibuka di Chrome: {link[:50]}...")
                self.close_tabs_button.setEnabled(True)
                
            except Exception as e:
                # Jika masih error, coba buat driver baru
                print(f"DEBUG: Error with existing driver, creating new one: {e}")
                try:
                    self.setup_chrome_driver_for_single_link()
                    self.chrome_driver.get(link)
                    tab_handle = self.chrome_driver.current_window_handle
                    self.opened_chrome_tabs.append(tab_handle)
                    
                    if link_index >= 0:
                        self.mark_link_opened(link_index)
                    
                    self.status_label.setText(f"Link dibuka di Chrome: {link[:50]}...")
                    self.close_tabs_button.setEnabled(True)
                    
                except Exception as e2:
                    QMessageBox.critical(self, "Error", f"Gagal buka link: {str(e2)}")
                    print(f"DEBUG: Error opening single link {link}: {e2}")
                
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Error saat buka link: {str(e)}")
            print(f"DEBUG: Error in open_single_link: {e}")
    
    def is_chrome_driver_valid(self):
        """Cek apakah Chrome driver session masih valid"""
        if not self.chrome_driver:
            return False
        
        try:
            # Coba akses session ID untuk test validitas
            _ = self.chrome_driver.session_id
            # Coba command sederhana untuk test koneksi
            _ = self.chrome_driver.current_url
            return True
        except Exception as e:
            print(f"DEBUG: Chrome driver session invalid: {e}")
            # Cleanup driver yang tidak valid
            try:
                self.chrome_driver.quit()
            except:
                pass
            self.chrome_driver = None
            return False
    
    def setup_chrome_driver_for_single_link(self):
        """Setup Chrome driver khusus untuk membuka link individual"""
        try:
            # Path ke chromedriver.exe di folder yang sama
            chromedriver_path = Path(__file__).parent / "chromedriver.exe"
            
            if not chromedriver_path.exists():
                raise Exception("chromedriver.exe tidak ditemukan di folder aplikasi")
            
            # Chrome options untuk incognito mode
            chrome_options = Options()
            chrome_options.add_argument("--incognito")
            chrome_options.add_argument("--disable-web-security")
            chrome_options.add_argument("--disable-features=VizDisplayCompositor")
            chrome_options.add_argument("--disable-blink-features=AutomationControlled")
            chrome_options.add_experimental_option("excludeSwitches", ["enable-automation"])
            chrome_options.add_experimental_option('useAutomationExtension', False)
            
            # Setup service dan driver
            service = Service(str(chromedriver_path))
            self.chrome_driver = webdriver.Chrome(service=service, options=chrome_options)
            self.chrome_driver.execute_script("Object.defineProperty(navigator, 'webdriver', {get: () => undefined})")
            
            # Reset daftar tab karena ini driver baru
            self.opened_chrome_tabs = []
            
            print(f"DEBUG: New Chrome driver created for single link")
            
        except Exception as e:
            print(f"DEBUG: Failed to setup Chrome driver for single link: {e}")
            raise e

def main():
    app = QApplication(sys.argv)
    
    # Set aplikasi ID untuk Windows taskbar agar ikon muncul dengan benar
    try:
        ctypes.windll.shell32.SetCurrentProcessExplicitAppUserModelID('LinkOpener.App.1.0')
    except:
        pass  # Ignore jika tidak bisa set app ID
    
    # Set aplikasi icon untuk semua dialog
    icon_path = Path(__file__).parent / "link_opener.ico" 
    if icon_path.exists():
        app.setWindowIcon(QIcon(str(icon_path)))
    
    # Set style
    app.setStyle('Fusion')
    
    window = LinkOpenerApp()
    window.show()
    
    sys.exit(app.exec())


if __name__ == "__main__":
    main()
